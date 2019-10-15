import os
from pptx import Presentation
from pptx.util import Inches
import re


def natural_keys(text):
    return [atoi(c) for c in re.split(r'(\d+)', text)]


def make_pptx(path):
    prs = Presentation()

    files = os.listdir(path)
    files = sorted(files, key=natural_keys)

    for f in files:
        print(f)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(0)
        slide.shapes.add_picture(path + '\\' + f, left, top, width=Inches(10), height=Inches(8.5))

    prs.save(path + '\\' + 'SeleniumSlide.pptx')
